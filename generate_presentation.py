"""
GRC Skills Benchmark — Executive PowerPoint Presentation
Generates a professional .pptx for director/partner audience.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# Colors
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x4E, 0x9A, 0xF5)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0x8B)
ACCENT_RED = RGBColor(0xE8, 0x5D, 0x5D)
ACCENT_ORANGE = RGBColor(0xF5, 0xA6, 0x23)
ACCENT_PURPLE = RGBColor(0xA8, 0x7C, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
VERY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
BG_LIGHT = RGBColor(0xFA, 0xFA, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_shape_box(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=DARK_GRAY, spacing=1.2, font_name="Calibri"):
    """Add textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, fsize, fcolor) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fsize if fsize else font_size)
        p.font.bold = bold
        p.font.color.rgb = fcolor if fcolor else color
        p.font.name = font_name
        p.space_after = Pt(font_size * spacing * 0.3)
    return txBox

# ============================================================
# SLIDE 1 — TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 1, 1.5, 11, 1.2, "GRC Skills para IA", 44, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1, 2.7, 11, 0.8, "Resultados do Benchmark de Performance", 28, False, ACCENT_BLUE, PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.0, 11, 0.6, "ISO 27001  |  GDPR  |  NIST CSF 2.0", 20, False, LIGHT_GRAY, PP_ALIGN.CENTER)
add_textbox(slide, 1, 5.5, 11, 0.5, "Abril 2026  |  54 avaliacoes  |  3 frameworks  |  Claude Opus 4.6", 14, False, MED_GRAY, PP_ALIGN.CENTER)
add_textbox(slide, 1, 6.3, 11, 0.4, "Fernando P. Marciano", 16, False, LIGHT_GRAY, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — O QUE SAO GRC SKILLS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "O que sao GRC Skills?", 32, True, DARK_GRAY, PP_ALIGN.LEFT)

# Divider line
from pptx.enum.shapes import MSO_SHAPE
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, 0.8, 1.4, 7, 1.0,
    "Skills sao modulos de conhecimento especializado que podem ser carregados em assistentes de IA para melhorar suas respostas em dominios especificos de compliance.",
    18, False, MED_GRAY)

# 3 cards
cards = [
    ("ISO 27001", "Gestao de Seguranca\nda Informacao", "93 controles\n4 temas", ACCENT_BLUE),
    ("GDPR", "Protecao de Dados\nPessoais (UE)", "99 artigos\n173 recitais", ACCENT_GREEN),
    ("NIST CSF 2.0", "Framework de\nCyberseguranca", "6 funcoes\n106 subcategorias", ACCENT_PURPLE),
]

for i, (title, desc, detail, color) in enumerate(cards):
    x = 0.8 + i * 4.0
    box = add_shape_box(slide, x, 2.8, 3.5, 3.5, VERY_LIGHT)
    # top color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.8), Inches(3.5), Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, x + 0.3, 3.1, 3, 0.5, title, 22, True, DARK_GRAY)
    add_textbox(slide, x + 0.3, 3.7, 3, 1.0, desc, 16, False, MED_GRAY)
    add_textbox(slide, x + 0.3, 4.8, 3, 0.8, detail, 14, False, color)

add_textbox(slide, 0.8, 6.6, 11, 0.5,
    "Pergunta central: A IA responde melhor sobre compliance quando tem uma skill carregada?",
    18, True, ACCENT_BLUE)

# ============================================================
# SLIDE 3 — METODOLOGIA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Como testamos", 32, True, DARK_GRAY)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Formula visual
formula_items = [
    ("9", "cenarios reais\nde compliance"),
    ("x 2", "condicoes\n(com e sem skill)"),
    ("x 3", "repeticoes\ncada"),
    ("= 54", "avaliacoes\nindependentes"),
]
for i, (num, desc) in enumerate(formula_items):
    x = 1.0 + i * 3.0
    add_textbox(slide, x, 1.5, 1.5, 0.8, num, 40, True, ACCENT_BLUE, PP_ALIGN.CENTER)
    add_textbox(slide, x - 0.3, 2.4, 2.0, 0.8, desc, 14, False, MED_GRAY, PP_ALIGN.CENTER)

# 5 criteria
add_textbox(slide, 0.8, 3.5, 11, 0.5, "5 criterios de qualidade (escala 1-5):", 18, True, DARK_GRAY)

criteria = [
    ("Precisao", "Informacoes corretas?", ACCENT_BLUE),
    ("Especificidade", "Cita controles exatos?", ACCENT_GREEN),
    ("Estrutura", "Formato auditavel?", ACCENT_PURPLE),
    ("Aplicabilidade", "Pratico e implementavel?", ACCENT_ORANGE),
    ("Completude", "Cobre tudo?", ACCENT_RED),
]
for i, (name, desc, color) in enumerate(criteria):
    x = 0.8 + i * 2.4
    box = add_shape_box(slide, x, 4.2, 2.1, 1.8, VERY_LIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.2), Inches(2.1), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, x + 0.15, 4.45, 1.9, 0.5, name, 16, True, DARK_GRAY, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 5.0, 1.9, 0.6, desc, 13, False, MED_GRAY, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4 — RESULTADO PRINCIPAL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Resultado Principal", 32, True, WHITE)

# Big number
add_textbox(slide, 1, 1.5, 11, 2.0, "+68%", 120, True, ACCENT_GREEN, PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.4, 11, 0.6, "melhoria geral na qualidade das respostas", 24, False, LIGHT_GRAY, PP_ALIGN.CENTER)

# Before/After boxes
box_before = add_shape_box(slide, 2, 4.5, 4, 2.0, RGBColor(0x2A, 0x2A, 0x45))
add_textbox(slide, 2.3, 4.6, 3.5, 0.4, "SEM skill", 16, True, ACCENT_RED)
add_textbox(slide, 2.3, 5.1, 3.5, 0.7, "14.6 / 25", 36, True, LIGHT_GRAY)
add_textbox(slide, 2.3, 5.8, 3.5, 0.4, "58% da pontuacao maxima", 14, False, MED_GRAY)

box_after = add_shape_box(slide, 7, 4.5, 4, 2.0, RGBColor(0x2A, 0x2A, 0x45))
add_textbox(slide, 7.3, 4.6, 3.5, 0.4, "COM skill", 16, True, ACCENT_GREEN)
add_textbox(slide, 7.3, 5.1, 3.5, 0.7, "24.6 / 25", 36, True, WHITE)
add_textbox(slide, 7.3, 5.8, 3.5, 0.4, "98% da pontuacao maxima", 14, False, ACCENT_GREEN)

# Arrow
add_textbox(slide, 5.8, 5.0, 1.5, 0.8, ">>>", 32, True, ACCENT_GREEN, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5 — MELHORIA POR CRITERIO (CHART)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Melhoria por Criterio", 32, True, DARK_GRAY)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Chart
chart_data = CategoryChartData()
chart_data.categories = ['Especificidade', 'Estrutura', 'Completude', 'Aplicabilidade', 'Precisao']
chart_data.add_series('Sem Skill', (1.96, 2.85, 2.96, 3.22, 3.63))
chart_data.add_series('Com Skill', (4.96, 4.85, 5.00, 4.78, 5.00))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.2),
    chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(14)

plot = chart.plots[0]
plot.gap_width = 80

series_baseline = plot.series[0]
series_baseline.format.fill.solid()
series_baseline.format.fill.fore_color.rgb = ACCENT_RED

series_skill = plot.series[1]
series_skill.format.fill.solid()
series_skill.format.fill.fore_color.rgb = ACCENT_GREEN

cat_axis = chart.category_axis
cat_axis.tick_labels.font.size = Pt(13)
val_axis = chart.value_axis
val_axis.maximum_scale = 5.5
val_axis.minimum_scale = 0
val_axis.tick_labels.font.size = Pt(12)

# Lift annotations on the right
lifts = [
    ("Especificidade", "+153%", "Cita controles, artigos, IDs exatos"),
    ("Estrutura", "+70%", "Tabelas, gap analysis, formato auditavel"),
    ("Completude", "+69%", "Nenhum aspecto relevante omitido"),
    ("Aplicabilidade", "+48%", "Owners, prazos, fases 30/60/90 dias"),
    ("Precisao", "+38%", "Baseline ja razoavel; skill refina"),
]

for i, (name, lift, desc) in enumerate(lifts):
    y = 1.7 + i * 1.05
    add_textbox(slide, 8.8, y, 1.5, 0.4, lift, 22, True, ACCENT_GREEN)
    add_textbox(slide, 10.3, y + 0.05, 2.8, 0.4, desc, 12, False, MED_GRAY)

# ============================================================
# SLIDE 6 — EXEMPLO CONCRETO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Exemplo: Incidente de Credenciais via Slack", 28, True, DARK_GRAY)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Left: without skill
box_l = add_shape_box(slide, 0.8, 1.4, 5.7, 5.5, RGBColor(0xFF, 0xF0, 0xF0))
bar_l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.7), Pt(5))
bar_l.fill.solid()
bar_l.fill.fore_color.rgb = ACCENT_RED
bar_l.line.fill.background()

add_textbox(slide, 1.0, 1.6, 5.3, 0.4, "SEM skill — Score: 13/25", 18, True, ACCENT_RED)
add_multi_text(slide, 1.0, 2.2, 5.3, 4.5, [
    ('"Revogue as credenciais, rotacione senhas, e', False, 14, MED_GRAY),
    ('treine a equipe sobre boas praticas de seguranca', False, 14, MED_GRAY),
    ('da informacao. Isso se relaciona com controles', False, 14, MED_GRAY),
    ('de acesso do ISO 27001."', False, 14, MED_GRAY),
    ('', False, 10, MED_GRAY),
    ('Problemas:', True, 15, ACCENT_RED),
    ('  - Nenhum controle especifico citado', False, 13, MED_GRAY),
    ('  - Sem formato de auditoria', False, 13, MED_GRAY),
    ('  - Sem timeline de acoes', False, 13, MED_GRAY),
    ('  - Omite DLP, monitoramento, incidentes', False, 13, MED_GRAY),
])

# Right: with skill
box_r = add_shape_box(slide, 6.8, 1.4, 5.7, 5.5, RGBColor(0xF0, 0xFF, 0xF0))
bar_r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.7), Pt(5))
bar_r.fill.solid()
bar_r.fill.fore_color.rgb = ACCENT_GREEN
bar_r.line.fill.background()

add_textbox(slide, 7.0, 1.6, 5.3, 0.4, "COM skill — Score: 25/25", 18, True, ACCENT_GREEN)
add_multi_text(slide, 7.0, 2.2, 5.3, 4.5, [
    ('Controles ISO 27001:2022 identificados:', True, 14, DARK_GRAY),
    ('  A.5.10 — Uso aceitavel de informacao', False, 12, MED_GRAY),
    ('  A.5.14 — Transferencia de informacao', False, 12, MED_GRAY),
    ('  A.5.17 — Informacao de autenticacao', False, 12, MED_GRAY),
    ('  A.8.12 — Prevencao vazamento (novo 2022)', False, 12, ACCENT_GREEN),
    ('  A.8.16 — Atividades de monitoramento', False, 12, MED_GRAY),
    ('  A.5.24 — Gestao de incidentes', False, 12, MED_GRAY),
    ('', False, 8, MED_GRAY),
    ('Acoes imediatas (0-24h): revogar, auditar Slack', False, 12, MED_GRAY),
    ('Curto prazo (1-2 sem): secrets manager, DLP', False, 12, MED_GRAY),
    ('Evidencias: logs Slack Admin, registro incidente', False, 12, MED_GRAY),
])

# ============================================================
# SLIDE 7 — RESULTADOS POR FRAMEWORK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Resultados por Framework", 32, True, DARK_GRAY)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Chart
chart_data2 = CategoryChartData()
chart_data2.categories = ['ISO 27001', 'GDPR', 'NIST CSF 2.0']
chart_data2.add_series('Sem Skill', (13.56, 15.22, 15.11))
chart_data2.add_series('Com Skill', (24.67, 24.56, 24.56))

chart2 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8), Inches(1.5), Inches(7), Inches(5.0),
    chart_data2
).chart

chart2.has_legend = True
chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
chart2.legend.font.size = Pt(14)

plot2 = chart2.plots[0]
plot2.gap_width = 100

s1 = plot2.series[0]
s1.format.fill.solid()
s1.format.fill.fore_color.rgb = ACCENT_RED

s2 = plot2.series[1]
s2.format.fill.solid()
s2.format.fill.fore_color.rgb = ACCENT_GREEN

chart2.value_axis.maximum_scale = 27
chart2.value_axis.minimum_scale = 0
chart2.value_axis.tick_labels.font.size = Pt(12)
chart2.category_axis.tick_labels.font.size = Pt(14)

# Callouts right side
fw_data = [
    ("ISO 27001", "+82%", "Maior ganho. Sistema de numeracao\ncomplexo (93 controles, 4 temas,\n11 novos em 2022) se beneficia\nmuito de dados de referencia.", ACCENT_BLUE),
    ("GDPR", "+61%", "Citacao obrigatoria de artigos\ntransforma respostas vagas em\norientacao juridica precisa.", ACCENT_GREEN),
    ("NIST CSF", "+62%", "Subcategorias, tiers, mapeamentos\ncruzados (800-53, CIS) agregam\nvalor significativo.", ACCENT_PURPLE),
]
for i, (name, lift, desc, color) in enumerate(fw_data):
    y = 1.8 + i * 1.8
    add_textbox(slide, 8.2, y, 1.5, 0.5, lift, 28, True, color)
    add_textbox(slide, 9.8, y + 0.1, 3.3, 1.4, desc, 12, False, MED_GRAY)

# ============================================================
# SLIDE 8 — CONFIABILIDADE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Confiabilidade Estatistica", 32, True, DARK_GRAY)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

stats = [
    ("54", "avaliacoes\nindependentes", "27 baseline + 27 com skill", ACCENT_BLUE),
    ("3x", "repeticoes\npor teste", "Media, desvio padrao, IC 95%", ACCENT_GREEN),
    ("0%", "overlap nos\nintervalos", "Diferenca significativa\nem TODOS os criterios", ACCENT_PURPLE),
    ("-46%", "reducao de\nvariancia", "Respostas mais\nprevisiveis com skill", ACCENT_ORANGE),
]

for i, (num, label, desc, color) in enumerate(stats):
    x = 0.8 + i * 3.1
    box = add_shape_box(slide, x, 1.5, 2.7, 3.5, VERY_LIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(2.7), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, x, 1.8, 2.7, 0.8, num, 44, True, color, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 2.7, 2.3, 0.8, label, 16, True, DARK_GRAY, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 3.6, 2.3, 1.0, desc, 13, False, MED_GRAY, PP_ALIGN.CENTER)

# IC 95% table visual
add_textbox(slide, 0.8, 5.3, 12, 0.5, "Intervalos de Confianca 95% — nenhuma sobreposicao entre condicoes", 16, True, DARK_GRAY)

ic_data = [
    ("Criterio", "Baseline IC 95%", "Com Skill IC 95%", "Significativo?"),
    ("Precisao", "[3.44 — 3.82]", "[5.00 — 5.00]", "Sim"),
    ("Especificidade", "[1.89 — 2.04]", "[4.89 — 5.04]", "Sim"),
    ("Estrutura", "[2.68 — 3.03]", "[4.72 — 4.99]", "Sim"),
    ("Aplicabilidade", "[3.06 — 3.38]", "[4.62 — 4.94]", "Sim"),
    ("Completude", "[2.89 — 3.04]", "[5.00 — 5.00]", "Sim"),
]

from pptx.util import Inches, Pt
tbl = slide.shapes.add_table(len(ic_data), 4, Inches(0.8), Inches(5.8), Inches(11), Inches(1.5)).table
tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(3.0)
tbl.columns[2].width = Inches(3.0)
tbl.columns[3].width = Inches(2.5)

for row_idx, row_data in enumerate(ic_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = tbl.cell(row_idx, col_idx)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = "Calibri"
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_BLUE
        elif col_idx == 3:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xFF, 0xE8)

# ============================================================
# SLIDE 9 — IMPLICACOES PARA O NEGOCIO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Implicacoes para o Negocio", 32, True, DARK_GRAY)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

implications = [
    ("Produtividade em Compliance",
     "Tarefas que exigiriam um consultor senior para formatar corretamente (controles, artigos, tabelas de gap analysis) podem ser feitas em minutos. Output ja sai no formato auditavel.",
     ACCENT_BLUE),
    ("Reducao de Risco de Omissao",
     "Sem skill, a IA consistentemente omite elementos criticos: controles novos do ISO 2022, requisitos de DPIA do GDPR para dados de saude, funcao Govern do NIST CSF 2.0. Com skill, completude = 100%.",
     ACCENT_RED),
    ("Escalabilidade e Padronizacao",
     "Skills sao modulares e reutilizaveis. Qualquer membro da equipe obtem o mesmo nivel de qualidade, independente da sua experiencia em GRC. Variancia reduzida em 46%.",
     ACCENT_GREEN),
    ("Custo Zero de Integracao",
     "Skills sao open source (GitHub). Ganho de +68% na qualidade sem investimento em infraestrutura, licencas ou treinamento adicional.",
     ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(implications):
    y = 1.5 + i * 1.4
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Pt(5), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, 1.1, y, 4, 0.4, title, 18, True, DARK_GRAY)
    add_textbox(slide, 1.1, y + 0.45, 11, 0.7, desc, 14, False, MED_GRAY)

# ============================================================
# SLIDE 10 — PROXIMOS PASSOS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_textbox(slide, 0.8, 0.4, 11, 0.7, "Proximos Passos", 32, True, DARK_GRAY)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

steps = [
    ("1", "Integrar as 3 skills validadas no workflow de compliance da equipe", "Imediato", ACCENT_GREEN),
    ("2", "Adicionar skills para SOC 2 e HIPAA (disponiveis no mesmo repositorio)", "Curto prazo", ACCENT_BLUE),
    ("3", "Customizar skills com politicas internas e controles especificos da organizacao", "Medio prazo", ACCENT_PURPLE),
    ("4", "Reavaliar trimestralmente conforme frameworks sao atualizados", "Continuo", ACCENT_ORANGE),
]

for i, (num, desc, timeline, color) in enumerate(steps):
    y = 1.5 + i * 1.35
    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.1), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, 2.0, y, 8, 0.5, desc, 18, False, DARK_GRAY)
    add_textbox(slide, 2.0, y + 0.5, 3, 0.3, timeline, 13, True, color)

# ============================================================
# SLIDE 11 — CONCLUSAO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.8, 11.5, 0.7, "Conclusao", 32, True, WHITE)

add_textbox(slide, 1.5, 2.0, 10, 3.5,
    "Skills de GRC melhoram a qualidade das respostas de IA em compliance em 68%, com o maior impacto em especificidade (+153%) e completude (+69%), transformando outputs genericos em entregas prontas para auditoria.",
    26, False, LIGHT_GRAY, PP_ALIGN.CENTER)

# Key numbers row
key_nums = [
    ("+68%", "melhoria\ngeral", ACCENT_GREEN),
    ("+153%", "especificidade", ACCENT_BLUE),
    ("100%", "completude\ncom skill", ACCENT_PURPLE),
    ("p < 0.001", "significancia\nestatistica", ACCENT_ORANGE),
]

for i, (num, label, color) in enumerate(key_nums):
    x = 1.0 + i * 3.0
    add_textbox(slide, x, 4.5, 2.5, 0.8, num, 36, True, color, PP_ALIGN.CENTER)
    add_textbox(slide, x, 5.3, 2.5, 0.6, label, 14, False, LIGHT_GRAY, PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.3, 11, 0.5,
    "Benchmark: 54 avaliacoes | 3 frameworks | 3x repeticoes | IC 95% | Claude Opus 4.6 | Abril 2026",
    12, False, MED_GRAY, PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.7, 11, 0.4,
    "Fernando P. Marciano — Relatorio tecnico completo disponivel",
    14, False, LIGHT_GRAY, PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = "F:/Projetos/GRC-Skills/GRC_Skills_Benchmark_Executivo.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
